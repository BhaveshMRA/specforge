from fastapi import FastAPI, HTTPException, UploadFile, File, Depends
from fastapi.concurrency import run_in_threadpool
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, Response
from pydantic import BaseModel
import httpx
import urllib.request
import urllib.error
import json
import os
import io
import base64
import re
import ast
import shutil
import subprocess
import tempfile
import zipfile
import sys
import socket
import time
import atexit
from dotenv import load_dotenv

# ── Database ──────────────────────────────────────────────────────────────────
from sqlalchemy import create_engine, Column, String, Text, DateTime
from sqlalchemy.orm import declarative_base, sessionmaker
from datetime import datetime
from uuid import uuid4

load_dotenv()

DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///./specforge.db")
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()


class SavedArch(Base):
    __tablename__ = "saved_archs"
    id           = Column(String, primary_key=True, default=lambda: str(uuid4()))
    name         = Column(String, nullable=False)
    project_name = Column(String)
    created_at   = Column(DateTime, default=datetime.utcnow)
    arch_json    = Column(Text, nullable=False)


Base.metadata.create_all(bind=engine)

# ── App ───────────────────────────────────────────────────────────────────────

app = FastAPI(title="SpecForge API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://localhost:3000",
        "http://localhost:8002",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

OLLAMA_API_KEY = os.getenv("OLLAMA_API_KEY")

SYSTEM_PROMPT = """You are a senior software architect. When given even a vague or one-line product idea, you deeply reason about what a complete, production-ready version needs — including the full user journey from input to output.

CRITICAL RULES:
1. The FIRST layer must ALWAYS be "User Experience" (color: blue) with two components:
   - id: "msg_input", name: "Message Input", tech matching the actual UI (e.g. "Resume Upload Form" for ATS, "Chat UI" for chatbot), purpose: "User submits here"
   - id: "response_view", name: "Response Display", tech: "Streaming UI", purpose: "User sees output here"
   msg_input connects_to the first backend component. response_view connects_to nothing.

2. The LAST processing layer must have at least one component whose connects_to includes "response_view".

3. sample_query: generate a REALISTIC example of what a real user would actually input to THIS specific product. NOT a generic chatbot message. For an ATS: "Upload John Smith resume for Senior Engineer role". For a chatbot: "Explain transformer attention in simple terms". For a research tool: "Find recent papers on RAG evaluation". Match the domain exactly.

4. Number every key_flow step: "Step 1: User submits X → Step 2: ..." showing the complete round trip.

Return ONLY valid JSON:
{"project_name":"string","summary":"2-3 sentences","sample_query":"realistic user input for THIS specific product","layers":[{"id":"snake_case","name":"Layer Name","color":"blue|purple|teal|amber|coral|green|gray","description":"max 8 words","components":[{"id":"snake_case","name":"Component Name","tech":"specific tech","purpose":"max 6 words","connects_to":["exact_ids"]}]}],"key_flows":["Step 1: ... → Step 2: ... full round trip"]}

Rules: 4-6 layers. 1-4 components per layer. Specific tech only. Colors: blue=frontend/UI, green=API/backend, purple=agent/AI logic, amber=LLM/model services, teal=data/storage, coral=auth, gray=infra."""

REFINE_PROMPT = """You are a senior software architect conducting a design review. You will receive an existing architecture JSON and a set of requested changes from the engineer.

Your job:
1. Preserve the existing architecture structure as much as possible
2. Apply ONLY the requested changes — add layers, modify components, update connections, etc.
3. Return the complete updated architecture in the exact same JSON schema
4. Update key_flows and sample_query to reflect the changes
5. Keep the same project_name unless the change fundamentally changes the product

Return ONLY valid JSON in the same schema as the input. No explanation, no markdown."""

REASON_PROMPT = """You are a senior software architect conducting a design review. 
You will receive an existing architecture JSON, a chat history, and a prompt from the engineer.

Determine if the engineer is asking a QUESTION (e.g. "where is the database?", "what parameters are there?") or explicitly commanding a STRUCTURAL MODIFICATION (e.g. "add a cache layer", "swap MySQL for PostgreSQL", "add those to the input component").

CRITICAL RULE: Default to "chat" type UNLESS the engineer uses explicit imperative action words instructing you to modify the architecture (e.g., "add", "change", "remove", "update"). If the engineer is just making a statement, providing context, listing items, or asking a question, you MUST use "chat" type.

Return ONLY valid JSON in this exact schema:
{
  "type": "chat" | "architecture",
  "message": "Conversational, helpful response if type is 'chat', otherwise an empty string",
  "architecture": { /* the fully updated architecture JSON if type is 'architecture', otherwise null */ }
}

If the type is 'architecture', you must preserve the existing architecture structure as much as possible, applying ONLY the requested changes, and updating key_flows/sample_query appropriately. Do not output markdown, just the raw JSON object."""

CODE_REASON_PROMPT = """You are a senior full-stack engineer debugging a real running application you have direct access to — not a hypothetical one.

You will receive: the complete current source files, and (when the app is currently running via "Run App") live diagnostics gathered by actually driving a headless browser against it just now — console errors, the result of an automated fill-every-input-and-submit sweep, the visible page text afterward, and a screenshot. You also get chat history and a message from the engineer.

Determine if the engineer is asking a QUESTION (chat) or reporting/implying something that needs a CODE FIX.

CRITICAL RULES:
- Default to "chat" type UNLESS you have enough evidence (from the live diagnostics or the engineer's own description) to make a specific, confident code change. If the app isn't running and the engineer hasn't described the bug precisely enough to act on, ask for what you need instead of guessing.
- Fix the root cause, not the symptom. If diagnostics show something like "X.filter is not a function", check whether the backend response shape actually matches what the frontend code expects (e.g. `{tasks: [...]}` vs a bare array) rather than just adding a defensive `Array.isArray` guard.
- Never invent a fix for a problem the diagnostics/engineer didn't show you evidence of.
- You have full control over the sandbox: every backend/main.py or backend/main.js, PLUS every additional service under backend/<name>/main.py|main.js, is started automatically whenever the app runs (in the same sandbox, reachable from each other over localhost) -- you never need the engineer to run a shell command themselves, and must never say you can't execute something. If a service-to-service call fails (e.g. ECONNREFUSED to another internal service), diagnose it as a code issue: a missing/broken dependency crashing that service on startup, a port mismatch against rule 16's convention (gateway on 8000, additional services on 8001, 8002, ...), or a wrong path/filename the sandbox wouldn't recognize as an entrypoint -- and fix the code directly.

Return ONLY valid JSON in this exact schema:
{
  "type": "chat" | "fix",
  "message": "Your diagnosis, and/or a plain-English summary of what you changed and why",
  "files": [ {"path": "...", "content": "..."} ] | null
}

If type is "fix", "files" must contain ONLY the files you actually changed or added (never files you left untouched) -- they'll be merged into the existing project by path. Keeping this list minimal matters: every file you include gets fully regenerated, which is slow and risks introducing unrelated mistakes in files that were already working. Each included file's full new content must still hold to the same engineering bar as initial generation: no placeholders, matching routes/imports between frontend and backend, Tailwind polish preserved, CORS allowing all origins. Do not output markdown, just the raw JSON object."""

CODER_PROMPT = """You are a senior full-stack engineer. You receive a system architecture JSON (layers, components, tech choices, connections) produced by an Architect agent, and you generate a minimal, runnable scaffold that implements it.

CRITICAL RULES:
1. Map each component's "tech" field to a real, working implementation — no placeholders, no TODOs, no lorem ipsum.
2. Keep it lean on the backend: one clear entry point (e.g. backend/main.py), plus only the extra files strictly needed (models, one route file) — unless the architecture genuinely names a separate backend service (e.g. an API gateway plus a distinct intelligence/worker layer), in which case follow rule 16 instead of forcing everything into one process. On the frontend, decompose into real components instead of one giant App.jsx — separate files for the header, each form, each panel/dashboard section, and reusable pieces (e.g. a Card, a StatBadge) if the UI has repeated visual patterns. 10-20 files total is normal for a UI with more than one screen section.
3. Prefer boring, standard choices consistent with the component's stated tech (e.g. "PostgreSQL" -> SQLAlchemy models + schema; "React" -> a working App.jsx that calls the backend). If a tech name is vague, pick the most common concrete implementation.
4. Code must actually run together: consistent imports, matching route paths between frontend fetch calls and backend routes, a requirements.txt / package.json listing exactly the dependencies used.
5. Setup/run commands go in README.md (see rule 15) — not a separate field.
6. Backend must be directly runnable with no extra flags. If Python: backend/main.py ending with `if __name__ == "__main__": import uvicorn; uvicorn.run(app, host="0.0.0.0", port=8000)` — do not rely on a separate `uvicorn` command. If Node: backend/main.js listening on port 8000 when run via `node main.js` — do not rely on nodemon/ts-node or a separate dev command. CORS must allow all origins either way.
7. If the frontend is React, it MUST be a Vite app, not Create React App: package.json scripts are exactly {"dev":"vite","build":"vite build"}, deps include "vite" and "@vitejs/plugin-react" as devDependencies, entry point is frontend/index.html at the project root (not frontend/public/index.html) loading frontend/src/main.jsx, and `npm run build` outputs to frontend/dist. All backend calls from the frontend must use the full URL http://localhost:8000/... (no relative paths, no proxy).
8. UI must look like a real shipped product, not a prototype — aim for the visual bar of a funded SaaS product's dashboard, not a form demo. Load Tailwind via `<script src="https://cdn.tailwindcss.com"></script>` in frontend/index.html and style every component with Tailwind utility classes — no unstyled default form elements, no bare black-on-white HTML. Also load a real typeface: `<link>` Google Fonts (e.g. Inter or a domain-appropriate pairing) in index.html and set it as the base font, don't leave it on the browser default. Use: a real color palette (not default blue links), card layouts with borders/shadows/rounded corners, consistent spacing, proper typography hierarchy, hover/focus/disabled states with smooth `transition` on every interactive element, a loading state while a request is in flight, an empty state when a list has zero items, and inline error messages on failure (not alert()). If React: use "lucide-react" for icons (add it as a dependency) instead of emoji or raw SVG — real icons read as more finished. If the domain has summarizable numbers (totals, counts, balances), show them as a row of stat cards at the top of the page, not buried in a list.
9. Never emit LaTeX or math markup ($...$, \\rightarrow, \\times, etc.) anywhere in UI-facing text — browsers don't render it and it shows up as literal garbage. Use plain Unicode symbols instead (→, ×, ±) or plain words.
10. Match the completeness of the architecture, not the minimum to compile: real client-side form validation with visible error messages, handle the actual edge cases implied by the domain (e.g. empty inputs, zero/negative amounts, duplicate names, division remainders when splitting), and give every list/table a populated, realistic empty state — not just a console.log and a TODO.
11. When building a request body object, never use bare object-shorthand for a computed/renamed value (e.g. `{ ...data, split_with }` when the variable is actually named `splitWith`) — that's an undefined-variable crash. Always write the key explicitly: `{ ...data, split_with: splitWith }`. Error-state UI elements must use Tailwind's red palette (text-red-*, bg-red-*, or border-red-*) and only for genuine errors, never for other purposes.
12. Imports are per-file, not global: when the frontend is split into multiple component files, every icon, hook, or helper a file's JSX actually references MUST be imported at the top of that exact file — importing it in App.jsx does not make it available in components/Whatever.jsx. Before finishing each component file, re-check every JSX tag and function call in it against that file's own import list.
13. In backend/requirements.txt, never pin an exact version (==) for a Python package that has a compiled/Rust extension (pydantic, fastapi, uvicorn, sqlalchemy, pillow, cryptography, numpy, etc.) — an old exact pin may have no prebuilt wheel for whatever Python version actually runs it, forcing a from-source compile that fails. List these bare (no version) or with a minimum floor (>=), never ==, so pip picks whatever version has a prebuilt wheel.
14. Seed the app with realistic ground-truth sample data for the domain instead of shipping it empty — a seed script, fixture file, or in-memory sample dataset with real, concrete content (actual task rows, actual product records, actual sample documents with real paragraph text for a RAG/knowledge app, etc.), not lorem ipsum and not just an empty table waiting for the first user action. The goal is that running the app for the first time already shows a populated, demoable product. If the architecture implies an external LLM/vector-DB/API dependency that has no credentials available in this environment, still wire the real integration code, but have it fall back to serving the seeded ground-truth data/answers so the demo works standalone.
15. Always include a root-level README.md file: what the product does, the tech stack per layer, project structure, and the exact setup/run commands for backend and frontend (this replaces any separate run_instructions field — do not omit it).
16. If the architecture calls for more than one independently-running backend service (e.g. an API gateway that proxies to a separate intelligence/worker/ML service), put each additional service in its own subdirectory under backend/ named after the service (e.g. backend/logic/main.py or backend/logic/main.js), each with its own requirements.txt or package.json scoped to that subdirectory. The primary/gateway service is always backend/main.py or backend/main.js on port 8000 (per rule 6); every additional service must listen on 127.0.0.1 at a distinct fixed port starting from 8001 and going up (8001, 8002, ...), since all backend services run together in the same sandbox and reach each other over localhost — never assume an additional service is externally reachable or needs its own public port.

Return ONLY valid JSON:
{"files":[{"path":"backend/main.py","content":"full file contents"},{"path":"README.md","content":"full file contents"}]}

No explanation, no markdown fences, just the raw JSON object."""


# ── Pydantic models ───────────────────────────────────────────────────────────

class SpecRequest(BaseModel):
    spec: str


class RefineRequest(BaseModel):
    existing_arch: dict
    feedback: str
    chat_history: list = []


class CodeRequest(BaseModel):
    architecture: dict


class RunRequest(BaseModel):
    files: list


class CodeReasonRequest(BaseModel):
    files: list
    feedback: str
    chat_history: list = []
    images: list = []  # data: URLs (e.g. "data:image/png;base64,...") pasted/dropped by the engineer


class GithubPushRequest(BaseModel):
    files: list
    token: str
    repo: str  # "owner/repo"
    mode: str  # "main" | "branch"
    branch_name: str = ""
    commit_message: str = "Push from SpecForge"


class SaveRequest(BaseModel):
    name: str
    arch: dict


# ── Helpers ───────────────────────────────────────────────────────────────────

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def call_ollama(messages: list, _retries: int = 2, timeout: float = None) -> dict:
    """Call the Ollama cloud API using stdlib urllib. No timeout by default —
    the original generation flow lets the model finish no matter how long that
    takes. Callers sending a much heavier prompt (e.g. full source + a
    screenshot on every debug turn) should pass an explicit timeout so a slow
    or stuck request fails with a clear error instead of hanging the request
    forever."""
    if not OLLAMA_API_KEY:
        raise HTTPException(status_code=500, detail="OLLAMA_API_KEY not set in .env")

    payload = json.dumps(
        {"model": "gemma4:31b-cloud", "messages": messages, "temperature": 0.3}
    ).encode("utf-8")

    req = urllib.request.Request(
        "https://ollama.com/v1/chat/completions",
        data=payload,
        headers={
            "Authorization": f"Bearer {OLLAMA_API_KEY}",
            "Content-Type": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            body = resp.read().decode("utf-8")
    except urllib.error.HTTPError as e:
        detail = e.read().decode("utf-8", errors="replace")
        raise HTTPException(status_code=e.code, detail=f"Ollama API error: {detail}")
    except urllib.error.URLError as e:
        raise HTTPException(status_code=502, detail=f"Could not reach Ollama API: {e.reason}")
    except TimeoutError:
        raise HTTPException(status_code=504, detail=f"Ollama API did not respond within {timeout}s")

    data = json.loads(body)
    raw = data["choices"][0]["message"]["content"]
    cleaned = raw.replace("```json", "").replace("```", "").strip()
    # Sanitize invalid \u escapes (e.g. \u followed by non-hex or end of string)
    cleaned = re.sub(r'\\u(?![0-9a-fA-F]{4})', r'\\\\u', cleaned)
    try:
        return json.loads(cleaned, strict=False)
    except json.JSONDecodeError as e:
        # ponytail: the model occasionally emits structurally malformed JSON
        # (e.g. a stray unescaped quote in generated code). Retrying regenerates
        # the whole response rather than attempting to repair broken JSON.
        if _retries > 0:
            return call_ollama(messages, _retries=_retries - 1, timeout=timeout)
        raise HTTPException(status_code=500, detail=f"Failed to parse JSON after retries: {e}\n\nRaw text: {cleaned}")


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/health")
async def health():
    return {"status": "ok", "key_set": bool(OLLAMA_API_KEY)}


@app.post("/api/architect")
async def architect(req: SpecRequest):
    return await run_in_threadpool(call_ollama, [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user", "content": f"Architect this: {req.spec}"},
    ])


@app.post("/api/refine")
async def refine(req: RefineRequest):
    return await run_in_threadpool(call_ollama, [
        {"role": "system", "content": REFINE_PROMPT},
        {"role": "user", "content": (
            f"Existing architecture:\n{json.dumps(req.existing_arch, indent=2)}\n\n"
            f"Requested changes: {req.feedback}"
        )},
    ])


@app.post("/api/reason")
async def reason(req: RefineRequest):
    history_str = ""
    if req.chat_history:
        history_str = "Chat History:\n" + "\n".join([f"{msg['role']}: {msg['message']}" for msg in req.chat_history]) + "\n\n"

    return await run_in_threadpool(call_ollama, [
        {"role": "system", "content": REASON_PROMPT},
        {"role": "user", "content": (
            f"Existing architecture:\n{json.dumps(req.existing_arch, indent=2)}\n\n"
            f"{history_str}"
            f"Engineer's prompt: {req.feedback}"
        )},
    ])


def _validate_file(path: str, content: str) -> dict:
    """Deterministic syntax check only — catches 'won't parse', not 'is wrong'.
    JSX/TS are skipped: node's --check can't parse JSX/TS without a transformer,
    and adding babel/tsc just for validation isn't worth the dependency."""
    ext = os.path.splitext(path)[1].lower()

    if ext == ".py":
        try:
            ast.parse(content)
            return {"path": path, "status": "valid"}
        except SyntaxError as e:
            return {"path": path, "status": "invalid", "error": f"SyntaxError: {e.msg} at line {e.lineno}"}

    if ext == ".json":
        try:
            json.loads(content)
            return {"path": path, "status": "valid"}
        except json.JSONDecodeError as e:
            return {"path": path, "status": "invalid", "error": f"JSONDecodeError: {e}"}

    if ext == ".js" and shutil.which("node"):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".js", delete=False) as f:
            f.write(content)
            tmp_path = f.name
        try:
            result = subprocess.run(["node", "--check", tmp_path], capture_output=True, text=True, timeout=10)
            if result.returncode == 0:
                return {"path": path, "status": "valid"}
            return {"path": path, "status": "invalid", "error": result.stderr.strip()[:500]}
        finally:
            os.unlink(tmp_path)

    return {"path": path, "status": "skipped"}


def _detect_backend_entry(files: list) -> tuple:
    """Backend isn't always Python — the Coder can pick Node depending on the
    architecture's stated tech. Detect which entrypoint got generated."""
    paths = {f["path"] for f in files}
    if "backend/main.py" in paths:
        return "python", "main.py"
    if "backend/main.js" in paths:
        return "node", "main.js"
    return None, None


def _detect_all_backend_entries(files: list) -> list:
    """Some architectures split the backend into multiple independently-running
    services -- e.g. an API gateway plus a separate intelligence/worker service
    in its own subdirectory -- that call each other over localhost. Detects
    every one of them, not just the primary backend/main.py|main.js, so the
    sandbox can start all of them. Otherwise the second service never starts
    and the gateway's calls to it fail with ECONNREFUSED at runtime, even
    though the primary backend boots fine on its own.
    Returns [(runtime, dir_prefix, entry_filename), ...], dir_prefix like
    "backend" or "backend/logic"."""
    paths = {f["path"] for f in files}
    entries = []
    if "backend/main.py" in paths:
        entries.append(("python", "backend", "main.py"))
    elif "backend/main.js" in paths:
        entries.append(("node", "backend", "main.js"))
    for p in sorted(paths):
        if not p.startswith("backend/"):
            continue
        base = p.rsplit("/", 1)[-1]
        dir_prefix = p.rsplit("/", 1)[0]
        if base not in ("main.py", "main.js") or dir_prefix == "backend":
            continue
        runtime = "python" if base == "main.py" else "node"
        entries.append((runtime, dir_prefix, base))
    return entries


def _multi_backend_cmd(entries: list) -> str:
    """Builds a shell command that installs each service's own deps (scoped to
    its own directory, if it has a requirements.txt/package.json) and starts
    all of them as background jobs in the same container -- so a gateway's
    calls to another service over localhost actually land on something."""
    parts = []
    for runtime, dir_prefix, entry in entries:
        rel = dir_prefix[len("backend"):].lstrip("/")
        cd = f"cd /app/{rel}" if rel else "cd /app"
        if runtime == "python":
            parts.append(f"({cd} && (pip install --quiet --break-system-packages -r requirements.txt 2>&1 || true) && python {entry}) &")
        else:
            parts.append(f"({cd} && (npm install 2>&1 || true) && node {entry}) &")
    parts.append("wait")
    return "\n".join(parts)


def _image_for_entries(entries: list) -> str:
    runtimes = {e[0] for e in entries}
    if len(runtimes) > 1:
        _ensure_multi_image()
        return DOCKER_MULTI_IMAGE
    return _docker_image_for(next(iter(runtimes)))


def _mounts_for_entries(entries: list) -> list:
    runtimes = {e[0] for e in entries}
    mounts = []
    if "python" in runtimes:
        mounts.append((DOCKER_PIP_CACHE_VOLUME, "/root/.cache/pip"))
    if "node" in runtimes:
        mounts.append((DOCKER_NPM_CACHE_VOLUME, "/root/.npm"))
    return mounts


def _ports_for_entries(entries: list) -> list:
    """The primary/gateway service is always RUN_BACKEND_PORT; every
    additional service gets the next port up, matching CODER_PROMPT rule 16's
    convention (8001, 8002, ...) -- so readiness-checking can wait for all of
    them, not just the primary."""
    return [RUN_BACKEND_PORT + i for i in range(len(entries))]


# ── Docker sandbox ──────────────────────────────────────────────────────────
# Generated code (backend install/boot, frontend install/build) never runs
# directly on the host anymore -- it runs inside a resource-capped, --rm
# Docker container with only its own temp dir bind-mounted in. Named volumes
# cache pip/npm downloads across runs (otherwise every install would be cold,
# which would defeat the point of caring about generation speed at all).
# ponytail: still no network isolation (pip/npm need the internet to fetch
# packages, so a malicious dependency could still exfiltrate at install time)
# -- that needs a private package mirror or allowlist, real infra beyond what
# a local single-user tool needs today.

DOCKER_LABEL = "specforge=1"
DOCKER_PYTHON_IMAGE = "python:3.13-slim"
DOCKER_NODE_IMAGE = "node:20-slim"
DOCKER_PIP_CACHE_VOLUME = "specforge-pip-cache"
DOCKER_NPM_CACHE_VOLUME = "specforge-npm-cache"
DOCKER_MULTI_IMAGE = "specforge-multi:latest"  # node:20-slim + python3 -- for backends split across runtimes
DOCKER_MULTI_DOCKERFILE = (
    "FROM node:20-slim\n"
    "RUN apt-get update && apt-get install -y --no-install-recommends python3 python3-pip "
    "&& ln -sf /usr/bin/python3 /usr/bin/python && rm -rf /var/lib/apt/lists/*\n"
)


def _ensure_multi_image():
    """Builds the combined Python+Node image once -- a no-op every time after,
    since Docker just finds the already-built image locally by tag."""
    exists = subprocess.run(["docker", "image", "inspect", DOCKER_MULTI_IMAGE], capture_output=True).returncode == 0
    if exists:
        return
    subprocess.run(
        ["docker", "build", "-t", DOCKER_MULTI_IMAGE, "-"],
        input=DOCKER_MULTI_DOCKERFILE, capture_output=True, text=True, timeout=180,
    )


def _docker_available() -> bool:
    if not shutil.which("docker"):
        return False
    try:
        return subprocess.run(["docker", "info"], capture_output=True, timeout=5).returncode == 0
    except Exception:
        return False


def _docker_image_for(runtime: str) -> str:
    return DOCKER_PYTHON_IMAGE if runtime == "python" else DOCKER_NODE_IMAGE


def _docker_reap_orphans():
    """Safety net for the exact bug we hit in dev: if this process restarts
    (e.g. uvicorn --reload) while a container is live, in-memory tracking is
    wiped before cleanup can run. Containers are labeled, so they can always
    be found and removed regardless of what our own state remembers."""
    if not _docker_available():
        return
    ids = subprocess.run(["docker", "ps", "-aq", "--filter", f"label={DOCKER_LABEL}"], capture_output=True, text=True)
    for cid in ids.stdout.split():
        subprocess.run(["docker", "rm", "-f", cid], capture_output=True)


def _docker_stop(name: str):
    subprocess.run(["docker", "rm", "-f", name], capture_output=True, timeout=15)


def _docker_run_sync(image: str, work_dir: str, cmd: str, mounts: list, timeout: int):
    """Runs `cmd` to completion inside a fresh container; work_dir is bind-mounted
    to /app. `mounts` is a list of (volume, container_path) cache mounts -- more
    than one when the container needs both a pip and an npm cache at once (a
    mixed-runtime multi-service backend). Returns a subprocess.CompletedProcess."""
    mount_args = []
    for volume, path in mounts:
        mount_args += ["-v", f"{volume}:{path}"]
    return subprocess.run(
        [
            "docker", "run", "--rm", "--label", DOCKER_LABEL,
            "-v", f"{work_dir}:/app", "-w", "/app",
            *mount_args,
            "--memory=1g", "--cpus=2",
            image, "sh", "-c", cmd,
        ],
        capture_output=True, text=True, timeout=timeout,
    )


def _docker_port_probe(image: str, port: int) -> list:
    """A liveness check run *inside* the container's own network namespace via
    `docker exec`. Checking the host-published port from outside is unreliable:
    Docker's port-forwarding proxy binds the host port as soon as the container
    starts, so a raw host-side connect() can succeed before (or even if) the
    app inside ever starts listening — it only proves the proxy is up, not the app."""
    if image == DOCKER_PYTHON_IMAGE:
        return ["python3", "-c", f"import socket; socket.create_connection(('127.0.0.1', {port}), timeout=1)"]
    return ["node", "-e", f"require('net').connect({port},'127.0.0.1').on('connect',()=>process.exit(0)).on('error',()=>process.exit(1))"]


def _docker_run_detached(image: str, work_dir: str, cmd: str, ports: list, name: str,
                          mounts: list, timeout: int) -> dict:
    """Starts `cmd` in a detached container publishing ports[0] (the primary
    gateway) to the host. `mounts` is a list of (volume, container_path) cache
    mounts. `ports` is every port a backend service in this container needs to
    be listening on -- for a multi-service backend, waits for ALL of them
    (each service installs its own deps in parallel, so the gateway's port
    opening first doesn't mean an internal service it calls is ready yet).
    Returns {"ok": True} once every port is accepting connections, else
    {"ok": False, "error": <logs>} with the container already cleaned up
    (including if it exits early, e.g. a crash on startup)."""
    _docker_stop(name)  # clear any stale container from a previous crashed run
    mount_args = []
    for volume, path in mounts:
        mount_args += ["-v", f"{volume}:{path}"]
    run = subprocess.run(
        [
            "docker", "run", "-d", "--label", DOCKER_LABEL, "--name", name,
            "-p", f"{ports[0]}:{ports[0]}",
            "-v", f"{work_dir}:/app", "-w", "/app",
            *mount_args,
            "--memory=1g", "--cpus=2",
            image, "sh", "-c", cmd,
        ],
        capture_output=True, text=True,
    )
    if run.returncode != 0:
        return {"ok": False, "error": f"docker run failed: {run.stderr.strip()[:800]}"}

    probes = [_docker_port_probe(image, p) for p in ports]
    deadline = time.time() + timeout
    while time.time() < deadline:
        running = subprocess.run(["docker", "inspect", "-f", "{{.State.Running}}", name], capture_output=True, text=True)
        if running.stdout.strip() != "true":
            break  # process inside exited (crashed) — stop polling, go straight to logs
        if all(subprocess.run(["docker", "exec", name] + probe, capture_output=True).returncode == 0 for probe in probes):
            return {"ok": True, "name": name}
        time.sleep(0.5)

    logs = subprocess.run(["docker", "logs", name], capture_output=True, text=True)
    _docker_stop(name)
    combined = (logs.stdout + logs.stderr).strip()
    return {"ok": False, "error": f"Container never opened port(s) {ports}: {combined[-1000:]}"}


def _boot_check_backend(files: list) -> dict:
    """Tier 2: actually try to start the generated backend (Python, Node, or
    both split across services) inside a sandboxed Docker container."""
    if not _docker_available():
        return {"status": "invalid", "error": "Docker is not running — start Docker Desktop to enable boot-checking"}

    backend_files = [f for f in files if f["path"].startswith("backend/")]
    entries = _detect_all_backend_entries(backend_files)
    if not entries:
        return {"status": "skipped"}

    with tempfile.TemporaryDirectory() as tmp:
        for f in backend_files:
            rel = f["path"][len("backend/"):]
            full = os.path.join(tmp, rel)
            os.makedirs(os.path.dirname(full) or tmp, exist_ok=True)
            with open(full, "w") as fh:
                fh.write(f["content"])

        cmd = _multi_backend_cmd(entries)
        result = _docker_run_detached(
            _image_for_entries(entries), tmp, cmd, _ports_for_entries(entries), "specforge-bootcheck",
            _mounts_for_entries(entries), timeout=150,
        )

        _docker_stop("specforge-bootcheck")
        if result["ok"]:
            return {"status": "valid"}
        return {"status": "invalid", "error": result["error"]}


def _run_frontend_smoke_test(frontend_url: str) -> dict:
    """Tier 4: fills every visible input/textarea with a dummy value, clicks
    the most likely submit button, and checks for (a) any console error or
    uncaught JS exception, or (b) a Tailwind red-* error element appearing
    afterward (relies on CODER_PROMPT rule 8's error-styling convention).
    ponytail: a heuristic smoke test, not a real E2E suite — catches 'the
    button doesn't actually work', not full user-flow coverage. Good enough
    for the failure mode we've actually seen (a silently-caught frontend bug
    that boot-checking the backend alone can't see)."""
    from playwright.sync_api import sync_playwright

    errors = []
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.on("console", lambda msg: errors.append(f"console.{msg.type}: {msg.text}") if msg.type == "error" else None)
            page.on("pageerror", lambda exc: errors.append(f"uncaught exception: {exc}"))

            page.goto(frontend_url, wait_until="networkidle", timeout=15000)

            for el in page.locator("input, textarea").all():
                try:
                    input_type = (el.get_attribute("type") or "text").lower()
                    if input_type in ("checkbox", "radio", "hidden", "file"):
                        continue
                    value = "1" if input_type == "number" else ("test@example.com" if input_type == "email" else "Test")
                    el.fill(value)
                except Exception:
                    pass  # not fillable (e.g. disabled) -- not itself a bug

            submit = page.locator(
                "button[type=submit], button:has-text('Add'), button:has-text('Submit'), "
                "button:has-text('Save'), button:has-text('Create')"
            ).first
            if submit.count() > 0:
                submit.click()
                try:
                    page.wait_for_load_state("networkidle", timeout=8000)
                except Exception:
                    pass  # request may still be genuinely slow, not necessarily a bug
                page.wait_for_timeout(500)  # let React finish re-rendering after the response lands

            red_error = page.locator("[class*='text-red-'], [class*='bg-red-'], [class*='border-red-']").first
            if red_error.count() > 0 and red_error.is_visible():
                errors.append(f"error state shown after submit: {red_error.inner_text()[:200]}")

            browser.close()
    except Exception as e:
        errors.append(f"smoke test crashed: {e}")

    if errors:
        return {"status": "invalid", "error": "; ".join(errors)[:800]}
    return {"status": "valid"}


def _capture_live_diagnostics(frontend_url: str) -> dict:
    """Drives a real headless browser against the app's own live preview window
    (the same one shown in the iframe) to gather ground truth for the debug chat
    -- this is the "browser control" the debug agent runs for itself before every
    turn, instead of asking the engineer to paste console output. Same
    fill-every-input-and-submit sweep as the Tier-4 smoke test (most runtime
    bugs only show up after an interaction, not on initial page load), plus a
    screenshot so a vision-capable model can see what the engineer sees."""
    from playwright.sync_api import sync_playwright

    console_errors = []
    screenshot_b64 = None
    page_text = ""
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.on("console", lambda msg: console_errors.append(f"console.{msg.type}: {msg.text}") if msg.type == "error" else None)
            page.on("pageerror", lambda exc: console_errors.append(f"uncaught exception: {exc}"))

            page.goto(frontend_url, wait_until="networkidle", timeout=15000)

            for el in page.locator("input, textarea").all():
                try:
                    input_type = (el.get_attribute("type") or "text").lower()
                    if input_type in ("checkbox", "radio", "hidden", "file"):
                        continue
                    value = "1" if input_type == "number" else ("test@example.com" if input_type == "email" else "Test")
                    el.fill(value)
                except Exception:
                    pass

            submit = page.locator(
                "button[type=submit], button:has-text('Add'), button:has-text('Submit'), "
                "button:has-text('Save'), button:has-text('Create')"
            ).first
            if submit.count() > 0:
                submit.click()
                try:
                    page.wait_for_load_state("networkidle", timeout=8000)
                except Exception:
                    pass
                page.wait_for_timeout(500)

            screenshot_b64 = base64.b64encode(page.screenshot()).decode()
            page_text = page.inner_text("body")[:3000]
            browser.close()
    except Exception as e:
        console_errors.append(f"diagnostic sweep crashed: {e}")

    return {"console_errors": console_errors, "screenshot_b64": screenshot_b64, "page_text": page_text}


def _run_checks(files: list) -> tuple:
    validations = [_validate_file(f["path"], f["content"]) for f in files]
    if any(v["status"] == "invalid" for v in validations):
        skipped = {"status": "skipped"}
        return validations, skipped, skipped, [v for v in validations if v["status"] == "invalid"]

    has_frontend = any(f["path"].startswith("frontend/") for f in files)
    frontend_check = {"status": "skipped"}

    if CURRENT_RUN is not None:
        # An app is live via Run App on these same ports -- never tear it
        # down just to run a background check. Skip rather than surprise-kill
        # the user's active session.
        boot_check = {"status": "skipped"}
    elif has_frontend:
        started = _start_full_stack(files)
        if not started["ok"]:
            boot_check = {"status": "invalid", "error": started["error"]}
        else:
            boot_check = {"status": "valid"}
            try:
                frontend_check = _run_frontend_smoke_test(started["frontend_url"])
            finally:
                _teardown_full_stack(started)
    else:
        boot_check = _boot_check_backend(files)

    failures = [v for v in validations if v["status"] == "invalid"]
    if boot_check["status"] == "invalid":
        failures.append({"path": "backend (boot)", "status": "invalid", "error": boot_check["error"]})
    if frontend_check["status"] == "invalid":
        failures.append({"path": "frontend (runtime)", "status": "invalid", "error": frontend_check["error"]})
    return validations, boot_check, frontend_check, failures


MAX_CODE_ATTEMPTS = 3


@app.post("/api/code")
async def generate_code(req: CodeRequest):
    messages = [
        {"role": "system", "content": CODER_PROMPT},
        {"role": "user", "content": (
            f"Architecture:\n{json.dumps(req.architecture, indent=2)}\n\nGenerate the code scaffold."
        )},
    ]

    result = await run_in_threadpool(call_ollama, messages)
    attempts = 1

    while True:
        validations, boot_check, frontend_check, failures = await run_in_threadpool(_run_checks, result.get("files", []))
        if not failures or attempts >= MAX_CODE_ATTEMPTS:
            break
        error_summary = "\n".join(f"- {v['path']}: {v['error']}" for v in failures)
        messages = messages + [
            {"role": "assistant", "content": json.dumps(result)},
            {"role": "user", "content": (
                f"These files failed validation:\n{error_summary}\n\n"
                "Return the complete corrected files JSON again in the same schema, fixing only these issues."
            )},
        ]
        result = await run_in_threadpool(call_ollama, messages)
        attempts += 1

    result["validation"] = validations
    result["boot_check"] = boot_check
    result["frontend_check"] = frontend_check
    result["attempts"] = attempts
    return result


# ── Live run (Tier 3): keep the generated app running and serve its UI ────────
# ponytail: one global run, not a per-user registry — this is a local single-
# user tool, not a hosted multi-tenant service. Backend runs sandboxed in
# Docker; the static file server for the already-built frontend stays a plain
# host subprocess since it only serves files, it doesn't execute generated code.

RUN_BACKEND_PORT = 8000
RUN_FRONTEND_PORT = 8001
RUN_BACKEND_CONTAINER = "specforge-run-backend"
CURRENT_RUN = None


def _wait_for_port(port: int, timeout: float) -> bool:
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            with socket.create_connection(("127.0.0.1", port), timeout=1):
                return True
        except OSError:
            time.sleep(0.3)
    return False


def _kill_port(port: int):
    """The static file server has no container label to reap by -- if this
    process restarts (e.g. uvicorn --reload) while it's live, the in-memory
    Popen handle is lost but the OS process survives, orphaned, still bound to
    the port. Worse, _wait_for_port then gives a false positive for the next
    run (something's listening -- just not the new server, which silently
    failed to bind and died). Kill whatever's actually on the port first."""
    try:
        pids = subprocess.run(["lsof", "-ti", f":{port}"], capture_output=True, text=True, timeout=5).stdout.split()
    except Exception:
        return
    for pid in pids:
        subprocess.run(["kill", "-9", pid], capture_output=True)


def _stop_current_run():
    global CURRENT_RUN
    if CURRENT_RUN is None:
        return
    _docker_stop(RUN_BACKEND_CONTAINER)
    static_proc = CURRENT_RUN.get("static_proc")
    if static_proc and static_proc.poll() is None:
        static_proc.terminate()
        try:
            static_proc.wait(timeout=5)
        except subprocess.TimeoutExpired:
            static_proc.kill()
    tmp_dir = CURRENT_RUN.get("tmp_dir")
    if tmp_dir:
        shutil.rmtree(tmp_dir, ignore_errors=True)
    CURRENT_RUN = None


_docker_reap_orphans()  # clean up anything orphaned by a previous crash/reload
_kill_port(RUN_FRONTEND_PORT)  # same, for the plain-subprocess static server
atexit.register(_stop_current_run)
atexit.register(_docker_reap_orphans)


def _start_full_stack(files: list) -> dict:
    """Boots backend (sandboxed in Docker) + builds/serves frontend. Caller
    owns teardown (via _teardown_full_stack) once done — this only starts
    things and never leaves a partial start running (cleans up on any failure)."""
    if not _docker_available():
        return {"ok": False, "error": "Docker is not running — start Docker Desktop to run generated apps"}

    tmp = tempfile.mkdtemp(prefix="specforge_run_")
    for f in files:
        full = os.path.join(tmp, f["path"])
        os.makedirs(os.path.dirname(full) or tmp, exist_ok=True)
        with open(full, "w") as fh:
            fh.write(f["content"])

    backend_dir = os.path.join(tmp, "backend")
    frontend_dir = os.path.join(tmp, "frontend")

    backend_files = [f for f in files if f["path"].startswith("backend/")]
    entries = _detect_all_backend_entries(backend_files)
    if not entries:
        shutil.rmtree(tmp, ignore_errors=True)
        return {"ok": False, "error": "No backend/main.py or backend/main.js in generated files"}
    if not os.path.isdir(frontend_dir):
        shutil.rmtree(tmp, ignore_errors=True)
        return {"ok": False, "error": "No frontend/ in generated files"}

    backend_cmd = _multi_backend_cmd(entries)
    backend_result = _docker_run_detached(
        _image_for_entries(entries), backend_dir, backend_cmd, _ports_for_entries(entries), RUN_BACKEND_CONTAINER,
        _mounts_for_entries(entries), timeout=180,
    )
    if not backend_result["ok"]:
        shutil.rmtree(tmp, ignore_errors=True)
        return {"ok": False, "error": f"Backend failed to start: {backend_result['error']}"}

    try:
        npm_install = _docker_run_sync(
            DOCKER_NODE_IMAGE, frontend_dir, "npm install",
            [(DOCKER_NPM_CACHE_VOLUME, "/root/.npm")], timeout=240,
        )
    except subprocess.TimeoutExpired:
        _docker_stop(RUN_BACKEND_CONTAINER)
        shutil.rmtree(tmp, ignore_errors=True)
        return {"ok": False, "error": "npm install timed out after 240s"}
    if npm_install.returncode != 0:
        _docker_stop(RUN_BACKEND_CONTAINER)
        shutil.rmtree(tmp, ignore_errors=True)
        return {"ok": False, "error": f"npm install failed: {npm_install.stderr.strip()[:800]}"}

    try:
        npm_build = _docker_run_sync(
            DOCKER_NODE_IMAGE, frontend_dir, "npm run build",
            [(DOCKER_NPM_CACHE_VOLUME, "/root/.npm")], timeout=180,
        )
    except subprocess.TimeoutExpired:
        _docker_stop(RUN_BACKEND_CONTAINER)
        shutil.rmtree(tmp, ignore_errors=True)
        return {"ok": False, "error": "npm run build timed out after 180s"}

    dist_dir = os.path.join(frontend_dir, "dist")
    if npm_build.returncode != 0 or not os.path.isdir(dist_dir):
        alt = os.path.join(frontend_dir, "build")  # in case the model ignored rule 7 and built a CRA app
        if os.path.isdir(alt):
            dist_dir = alt
        else:
            _docker_stop(RUN_BACKEND_CONTAINER)
            shutil.rmtree(tmp, ignore_errors=True)
            # esbuild/rollup put the actual "file:line: error" near the top of
            # stderr, not the tail (which is just the bundler's own stack trace)
            return {"ok": False, "error": f"npm run build failed: {npm_build.stderr.strip()[:1200]}"}

    _kill_port(RUN_FRONTEND_PORT)  # evict any orphan before binding, or _wait_for_port below would false-positive on it
    static_proc = subprocess.Popen(
        [sys.executable, "-m", "http.server", str(RUN_FRONTEND_PORT), "--directory", dist_dir],
        # ponytail: nothing ever reads these pipes -- http.server logs every
        # request to stderr, and an unread PIPE fills its OS buffer and
        # deadlocks the server thread mid-response after enough requests.
        # We don't need the output, so don't capture it at all.
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
    )
    if not _wait_for_port(RUN_FRONTEND_PORT, timeout=10):
        _docker_stop(RUN_BACKEND_CONTAINER)
        static_proc.kill()
        shutil.rmtree(tmp, ignore_errors=True)
        return {"ok": False, "error": "Static file server for the frontend failed to start"}

    return {
        "ok": True, "tmp": tmp, "static_proc": static_proc,
        "backend_url": f"http://localhost:{RUN_BACKEND_PORT}",
        "frontend_url": f"http://localhost:{RUN_FRONTEND_PORT}",
    }


def _teardown_full_stack(started: dict):
    _docker_stop(RUN_BACKEND_CONTAINER)
    static_proc = started.get("static_proc")
    if static_proc and static_proc.poll() is None:
        static_proc.terminate()
        try:
            static_proc.wait(timeout=5)
        except subprocess.TimeoutExpired:
            static_proc.kill()
    if started.get("tmp"):
        shutil.rmtree(started["tmp"], ignore_errors=True)


def _run_stack(files: list) -> dict:
    global CURRENT_RUN
    _stop_current_run()
    started = _start_full_stack(files)
    if not started["ok"]:
        return {"status": "error", "error": started["error"]}
    CURRENT_RUN = {"tmp_dir": started["tmp"], "static_proc": started["static_proc"]}
    return {"status": "running", "backend_url": started["backend_url"], "frontend_url": started["frontend_url"]}


@app.post("/api/code/run")
async def run_code(req: RunRequest):
    return await run_in_threadpool(_run_stack, req.files)


@app.post("/api/code/stop")
async def stop_code():
    await run_in_threadpool(_stop_current_run)
    return {"status": "stopped"}


@app.post("/api/code/reason")
async def code_reason(req: CodeReasonRequest):
    """The debug chat's engine: actually drives a headless browser against the
    live app before answering -- this is what gives the agent "direct access"
    instead of asking the engineer to describe/paste errors. It starts the app
    itself if it isn't already running (the engineer shouldn't have to click
    Run App first for the debugger to be able to see anything), and after
    applying any fix it (re)starts the app again so the result is immediately
    visible in the live preview, not just silently saved to the file list."""
    start_error = None
    if CURRENT_RUN is None:
        started = await run_in_threadpool(_run_stack, req.files)
        if started.get("status") != "running":
            start_error = started.get("error")

    diagnostics = None
    if CURRENT_RUN is not None:
        diagnostics = await run_in_threadpool(_capture_live_diagnostics, f"http://localhost:{RUN_FRONTEND_PORT}")

    history_str = ""
    if req.chat_history:
        history_str = "Chat History:\n" + "\n".join(f"{m['role']}: {m['message']}" for m in req.chat_history) + "\n\n"

    files_str = "\n\n".join(f"--- {f['path']} ---\n{f['content']}" for f in req.files)

    if diagnostics:
        diag_str = (
            "Live diagnostics just captured by driving a headless browser against the running app:\n"
            f"Console errors: {'; '.join(diagnostics['console_errors']) or '(none)'}\n\n"
            f"Visible page text after an automated fill-every-input-and-submit sweep:\n{diagnostics['page_text']}"
        )
    elif start_error:
        diag_str = f"Tried to start the app to inspect it live, but it failed to boot: {start_error}\n\nDiagnose from that boot failure, the source files, and the engineer's description."
    else:
        diag_str = "The app is not currently running, so no live diagnostics are available -- rely on the engineer's description and the source files."

    text_block = f"Current source files:\n{files_str}\n\n{diag_str}\n\n{history_str}Engineer's message: {req.feedback}"

    image_parts = []
    if diagnostics and diagnostics.get("screenshot_b64"):
        image_parts.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64,{diagnostics['screenshot_b64']}"}})
    for data_url in req.images:
        image_parts.append({"type": "image_url", "image_url": {"url": data_url}})

    if image_parts:
        note = " A screenshot of the running app after the interaction sweep is attached." if diagnostics and diagnostics.get("screenshot_b64") else ""
        if req.images:
            note += f" The engineer also attached {len(req.images)} image(s) — inspect them, they may show the bug directly."
        user_content = [{"type": "text", "text": text_block + note}] + image_parts
    else:
        user_content = text_block

    result = await run_in_threadpool(call_ollama, [
        {"role": "system", "content": CODE_REASON_PROMPT},
        {"role": "user", "content": user_content},
    ], timeout=240)

    if result.get("type") == "fix" and result.get("files"):
        merged = {f["path"]: f for f in req.files}
        for f in result["files"]:
            merged[f["path"]] = f
        result["files"] = list(merged.values())
        result["validation"] = [_validate_file(f["path"], f["content"]) for f in result["files"]]
        restart = await run_in_threadpool(_run_stack, result["files"])
        result["run_result"] = restart

    result["app_running"] = CURRENT_RUN is not None
    if CURRENT_RUN is not None:
        result["backend_url"] = f"http://localhost:{RUN_BACKEND_PORT}"
        result["frontend_url"] = f"http://localhost:{RUN_FRONTEND_PORT}"

    return result


@app.post("/api/code/zip")
async def zip_code(req: RunRequest):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for f in req.files:
            zf.writestr(f["path"], f["content"])
    return Response(
        content=buf.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": "attachment; filename=specforge-project.zip"},
    )


GITHUB_API = "https://api.github.com"


async def _github_request(client: httpx.AsyncClient, method: str, path: str, token: str, **kwargs) -> dict:
    r = await client.request(
        method, f"{GITHUB_API}{path}",
        headers={
            "Authorization": f"Bearer {token}",
            "Accept": "application/vnd.github+json",
            "X-GitHub-Api-Version": "2022-11-28",
        },
        **kwargs,
    )
    if r.status_code >= 400:
        try:
            detail = r.json().get("message", r.text)
        except Exception:
            detail = r.text
        raise HTTPException(status_code=r.status_code, detail=f"GitHub API error ({method} {path}): {detail}")
    return r.json() if r.content else {}


async def _upsert_ref(client: httpx.AsyncClient, owner: str, repo: str, token: str, branch: str, sha: str):
    """Points `branch` at `sha`, creating the branch if it doesn't exist yet
    (e.g. a brand-new branch, or the very first commit to an empty repo)."""
    try:
        await _github_request(client, "PATCH", f"/repos/{owner}/{repo}/git/refs/heads/{branch}", token, json={"sha": sha})
    except HTTPException as e:
        if e.status_code == 404:
            await _github_request(client, "POST", f"/repos/{owner}/{repo}/git/refs", token, json={"ref": f"refs/heads/{branch}", "sha": sha})
        else:
            raise


@app.post("/api/code/github-push")
async def github_push(req: GithubPushRequest):
    """Pushes the generated files as a single commit via the Git Data API
    (blobs/tree/commit), not one commit per file via the Contents API --
    layered on top of the target branch's existing tree so it doesn't clobber
    anything already in the repo that isn't part of this generation."""
    if "/" not in req.repo:
        raise HTTPException(status_code=400, detail="Repo must be in 'owner/repo' format")
    owner, repo = req.repo.split("/", 1)
    if req.mode not in ("main", "branch"):
        raise HTTPException(status_code=400, detail="mode must be 'main' or 'branch'")

    async with httpx.AsyncClient(timeout=30.0) as client:
        repo_info = await _github_request(client, "GET", f"/repos/{owner}/{repo}", req.token)
        default_branch = repo_info.get("default_branch", "main")
        target_branch = default_branch if req.mode == "main" else (req.branch_name.strip() or f"specforge-{int(time.time())}")

        # Base on the default branch's current tip -- absent entirely for a brand-new empty repo
        base_commit_sha = None
        base_tree_sha = None
        try:
            ref = await _github_request(client, "GET", f"/repos/{owner}/{repo}/git/ref/heads/{default_branch}", req.token)
            base_commit_sha = ref["object"]["sha"]
            base_commit = await _github_request(client, "GET", f"/repos/{owner}/{repo}/git/commits/{base_commit_sha}", req.token)
            base_tree_sha = base_commit["tree"]["sha"]
        except HTTPException as e:
            if e.status_code != 404:
                raise

        tree_body = {"tree": [{"path": f["path"], "mode": "100644", "type": "blob", "content": f["content"]} for f in req.files]}
        if base_tree_sha:
            tree_body["base_tree"] = base_tree_sha
        new_tree = await _github_request(client, "POST", f"/repos/{owner}/{repo}/git/trees", req.token, json=tree_body)

        commit_body = {"message": req.commit_message, "tree": new_tree["sha"]}
        if base_commit_sha:
            commit_body["parents"] = [base_commit_sha]
        new_commit = await _github_request(client, "POST", f"/repos/{owner}/{repo}/git/commits", req.token, json=commit_body)

        await _upsert_ref(client, owner, repo, req.token, target_branch, new_commit["sha"])

    return {
        "status": "pushed",
        "branch": target_branch,
        "commit_sha": new_commit["sha"],
        "commit_url": f"https://github.com/{owner}/{repo}/commit/{new_commit['sha']}",
        "branch_url": f"https://github.com/{owner}/{repo}/tree/{target_branch}",
    }


# ── File extraction helpers ───────────────────────────────────────────────────

def extract_pdf(content: bytes) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t.strip())
    return "\n\n".join(text_parts) or "[PDF had no extractable text]"


def extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs) or "[Document had no extractable text]"


def extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides) or "[Presentation had no extractable text]"


def extract_image_text(content: bytes, mime: str) -> str:
    """Send image to Gemma vision to describe/extract content."""
    b64 = base64.b64encode(content).decode()
    return f"data:{mime};base64,{b64}"  # returned to caller to embed in vision message


# ── Upload endpoint ────────────────────────────────────────────────────────────

ALLOWED_TYPES = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}


@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    ext = "." + file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
    if ext not in ALLOWED_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{ext}'. Allowed: PDF, DOCX, PPTX, PNG, JPG, WEBP"
        )

    content = await file.read()

    if ext == ".pdf":
        extracted = extract_pdf(content)
    elif ext == ".docx":
        extracted = extract_docx(content)
    elif ext == ".pptx":
        extracted = extract_pptx(content)
    elif ext in IMAGE_EXTS:
        # For images: send to Gemma vision to extract/describe content
        if not OLLAMA_API_KEY:
            raise HTTPException(status_code=500, detail="OLLAMA_API_KEY not set")
        b64 = base64.b64encode(content).decode()
        data_url = f"data:{ALLOWED_TYPES[ext]};base64,{b64}"
        async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
            r = await client.post(
                "https://ollama.com/v1/chat/completions",
                headers={"Authorization": f"Bearer {OLLAMA_API_KEY}",
                         "Content-Type": "application/json"},
                json={
                    "model": "gemma4:31b-cloud",
                    "messages": [{
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": data_url}},
                            {"type": "text", "text": (
                                "Extract and describe all text, diagrams, and relevant content "
                                "from this image. Be thorough — this will be used as context "
                                "for a software architecture tool."
                            )},
                        ],
                    }],
                    "temperature": 0.1,
                },
            )
        if r.status_code != 200:
            raise HTTPException(status_code=500, detail=f"Vision extraction failed: {r.text}")
        extracted = r.json()["choices"][0]["message"]["content"]
    else:
        extracted = "[File type not supported for extraction]"

    char_count = len(extracted)
    return {
        "filename": file.filename,
        "extracted_text": extracted,
        "char_count": char_count,
    }


# ── Save endpoints ─────────────────────────────────────────────────────────────

from fastapi import Depends
from sqlalchemy.orm import Session


@app.post("/api/saves", status_code=201)
def create_save(req: SaveRequest, db: Session = Depends(get_db)):
    save = SavedArch(
        id=str(uuid4()),
        name=req.name,
        project_name=req.arch.get("project_name"),
        arch_json=json.dumps(req.arch),
    )
    db.add(save)
    db.commit()
    db.refresh(save)
    return {"id": save.id, "name": save.name, "project_name": save.project_name,
            "created_at": save.created_at.isoformat()}


@app.get("/api/saves")
def list_saves(db: Session = Depends(get_db)):
    saves = db.query(SavedArch).order_by(SavedArch.created_at.desc()).all()
    return [{"id": s.id, "name": s.name, "project_name": s.project_name,
             "created_at": s.created_at.isoformat()} for s in saves]


@app.get("/api/saves/{save_id}")
def get_save(save_id: str, db: Session = Depends(get_db)):
    save = db.query(SavedArch).filter(SavedArch.id == save_id).first()
    if not save:
        raise HTTPException(status_code=404, detail="Save not found")
    return {"id": save.id, "name": save.name, "arch": json.loads(save.arch_json),
            "created_at": save.created_at.isoformat()}


@app.put("/api/saves/{save_id}")
def update_save(save_id: str, req: SaveRequest, db: Session = Depends(get_db)):
    save = db.query(SavedArch).filter(SavedArch.id == save_id).first()
    if not save:
        raise HTTPException(status_code=404, detail="Save not found")
    
    save.project_name = req.arch.get("project_name")
    save.arch_json = json.dumps(req.arch)
    db.commit()
    db.refresh(save)
    return {"id": save.id, "name": save.name, "project_name": save.project_name,
            "created_at": save.created_at.isoformat()}


@app.delete("/api/saves/{save_id}", status_code=204)
def delete_save(save_id: str, db: Session = Depends(get_db)):
    save = db.query(SavedArch).filter(SavedArch.id == save_id).first()
    if not save:
        raise HTTPException(status_code=404, detail="Save not found")
    db.delete(save)
    db.commit()


# ── Serve built React frontend ─────────────────────────────────────────────────

FRONTEND_DIST = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "frontend", "dist")

if os.path.isdir(FRONTEND_DIST):
    app.mount("/", StaticFiles(directory=FRONTEND_DIST, html=True), name="frontend")
